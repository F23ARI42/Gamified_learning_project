import os
import re
from flask import Flask, request, jsonify, render_template, redirect, url_for
from PyPDF2 import PdfReader
from pptx import Presentation
import google.generativeai as genai

# API Key configuration
my_api_key = "AIzaSyAszxQkapfFkay9FL0lOsKPGHAYc7_T5L4"
genai.configure(api_key=my_api_key)

# Flask app setup
app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Temporary storage for questions and answers
questions_and_answers = {}

# Function to extract text from PDF
def extract_text_from_pdf(pdf_file):
    try:
        reader = PdfReader(pdf_file)
        text = ''.join(page.extract_text() for page in reader.pages)
        return text.strip()
    except Exception as e:
        return f"Error extracting text from PDF: {e}"

# Function to extract text from PPT
def extract_text_from_ppt(ppt_file):
    try:
        prs = Presentation(ppt_file)
        text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame)
        return text.strip()
    except Exception as e:
        return f"Error extracting text from PPT: {e}"

# Function to generate questions
def generate_questions(text, difficulty, num_questions):
    prompt = {
        "Easy": f"Generate {num_questions} simple questions based on this text: {text}",
        "Medium": f"Generate {num_questions} moderately challenging questions based on this text: {text}",
        "Hard": f"Generate {num_questions} difficult questions based on this text: {text}"
    }.get(difficulty, "Invalid difficulty level. Choose Easy, Medium, or Hard.")

    try:
        model = genai.GenerativeModel('gemini-pro')
        response = model.generate_content(prompt)
        questions = re.sub(r'[*#]', '', response.text).strip().split("\n")
        return questions
    except Exception as e:
        return f"Error generating questions: {e}"

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate():
    try:
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file uploaded"}), 400

        difficulty = request.form.get('difficulty', 'Medium')
        num_questions = int(request.form.get('num_questions', 5))

        # Save and process file
        file_path = os.path.join(UPLOAD_FOLDER, file.filename)
        file.save(file_path)

        if file.filename.endswith('.pdf'):
            file_text = extract_text_from_pdf(file_path)
        elif file.filename.endswith('.pptx'):
            file_text = extract_text_from_ppt(file_path)
        else:
            return jsonify({"error": "Unsupported file type. Upload PDF or PPTX."}), 400

        if file_text.startswith("Error"):
            return jsonify({"error": file_text}), 500

        # Generate questions
        questions = generate_questions(file_text, difficulty, num_questions)

        if isinstance(questions, str) and questions.startswith("Error"):
            return jsonify({"error": questions}), 500

        # Save questions to the temporary storage
        questions_and_answers['questions'] = questions
        questions_and_answers['answers'] = [None] * len(questions)  # Placeholder for answers

        return redirect(url_for('questions_page'))

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/questions', methods=['GET', 'POST'])
def questions_page():
    if request.method == 'POST':
        # Capture user answers
        for idx, question in enumerate(questions_and_answers['questions']):
            questions_and_answers['answers'][idx] = request.form.get(f'answer_{idx}')
        return redirect(url_for('report_page'))
    return render_template('question.html', questions=questions_and_answers['questions'])

@app.route('/report', methods=['GET'])
def report_page():
    try:
        # Retrieve user answers and questions
        user_answers = questions_and_answers.get('answers', [])
        questions = questions_and_answers.get('questions', [])

        # Dynamically generate placeholder correct answers (replace with actual logic)
        correct_answers = [f"Answer{i+1}" for i in range(len(questions))]

        # Normalize answers for comparison
        def normalize(text):
            return text.strip().lower() if text else ""

        # Calculate the score and prepare results
        score = 0
        results = []

        for i, (question, user_answer) in enumerate(zip(questions, user_answers)):
            is_correct = normalize(user_answer) == normalize(correct_answers[i])
            results.append({
                "question": question,
                "user_answer": user_answer if user_answer else "None",
                "correct_answer": correct_answers[i],
                "is_correct": is_correct
            })
            if is_correct:
                score += 10  # Example reward: 10 points per correct answer

        return render_template('report.html', results=results, score=score)

    except Exception as e:
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    app.run(debug=True)

